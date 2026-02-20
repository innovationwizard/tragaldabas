"""Stage 7: Executive Output"""

from pathlib import Path
from typing import Dict, Any
from datetime import datetime

from core.interfaces import Stage
from core.models import AnalysisResult, OutputResult, NarrativeEvidence
from core.exceptions import StageError
from config import settings


class OutputManager(Stage[AnalysisResult, OutputResult]):
    """Stage 7: Generate executive output (text, markdown, PowerPoint)"""
    
    @property
    def name(self) -> str:
        return "Executive Output"
    
    @property
    def stage_number(self) -> int:
        return 7
    
    def validate_input(self, input_data: AnalysisResult) -> bool:
        return isinstance(input_data, AnalysisResult)
    
    async def execute(self, input_data: AnalysisResult) -> OutputResult:
        """Execute output generation"""
        output_dir = settings.get_output_path("insights")
        presentations_dir = settings.get_output_path("presentations")
        
        # Generate text output
        text_file = output_dir / "insights.txt"
        markdown_file = output_dir / "insights.md"
        
        self._generate_text_output(input_data, text_file)
        self._generate_markdown_output(input_data, markdown_file)
        
        # Generate PowerPoint
        pptx_file = presentations_dir / f"presentation_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
        self._generate_powerpoint(input_data, pptx_file)
        
        return OutputResult(
            text_file_path=str(text_file),
            markdown_file_path=str(markdown_file),
            pptx_file_path=str(pptx_file),
            slide_count=len(input_data.insights) + 2,  # Title + summary + insights
            insight_count=len(input_data.insights)
        )
    
    def _generate_text_output(self, analysis: AnalysisResult, file_path: Path):
        """Generate plain text output"""
        lines = [
            f"Tragaldabas Analysis Report",
            f"Domain: {analysis.domain.value}",
            f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}",
            "",
            "=" * 50,
            "EXECUTIVE SUMMARY",
            "=" * 50,
            ""
        ]

        # The Genius Move (Alpha Strike)
        if analysis.genius_insight and analysis.genius_insight.thesis:
            g = analysis.genius_insight
            lines.append("THE GENIUS MOVE (Strategic Alpha)")
            lines.append("-" * 40)
            lines.append(f"Thesis: {g.thesis}")
            if g.mechanism:
                lines.append(f"Mechanism: {g.mechanism}")
            if g.market_confluence:
                lines.append(f"Market Confluence: {g.market_confluence}")
            if g.estimated_upside:
                lines.append(f"Estimated Upside: {g.estimated_upside}")
            if g.kill_switch:
                lines.append(f"Kill Switch: {g.kill_switch}")
            lines.append("")
        
        # Top insights
        for i, insight in enumerate(analysis.insights[:5], 1):
            lines.append(f"{i}. {insight.headline}")
            lines.append(f"   {insight.detail}")
            lines.append(f"   Implication: {insight.implication}")
            lines.append("")
        
        lines.append("=" * 50)
        lines.append("DETAILED INSIGHTS")
        lines.append("=" * 50)
        lines.append("")
        
        for insight in analysis.insights:
            lines.append(f"• {insight.headline}")
            lines.append(f"  {insight.detail}")
            ev = insight.evidence
            if isinstance(ev, NarrativeEvidence):
                lines.append(f"  Evidence ({ev.source_type}): {ev.reference}")
                if ev.speaker:
                    lines.append(f"  Speaker: {ev.speaker}")
            else:
                lines.append(f"  Evidence: {ev.metric} = {ev.value}")
                if ev.delta is not None:
                    lines.append(f"  Change: {ev.delta_percent:.1f}%")
            lines.append(f"  Implication: {insight.implication}")
            lines.append("")
        
        with open(file_path, 'w', encoding='utf-8') as f:
            f.write("\n".join(lines))
    
    def _generate_markdown_output(self, analysis: AnalysisResult, file_path: Path):
        """Generate markdown output"""
        lines = [
            f"# Tragaldabas Analysis Report",
            "",
            f"**Domain:** {analysis.domain.value}",
            f"**Generated:** {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}",
            "",
            "## Executive Summary",
            ""
        ]

        # The Genius Move (Alpha Strike)
        if analysis.genius_insight and analysis.genius_insight.thesis:
            g = analysis.genius_insight
            lines.append("### ⚡ The Genius Move (Strategic Alpha)")
            lines.append("")
            lines.append(f"**Thesis:** {g.thesis}")
            if g.mechanism:
                lines.append(f"**Mechanism:** {g.mechanism}")
            if g.market_confluence:
                lines.append(f"**Market Confluence:** {g.market_confluence}")
            if g.estimated_upside:
                lines.append(f"**Estimated Upside:** {g.estimated_upside}")
            if g.kill_switch:
                lines.append(f"**Kill Switch:** {g.kill_switch}")
            lines.append("")
        
        for i, insight in enumerate(analysis.insights[:5], 1):
            lines.append(f"{i}. **{insight.headline}**")
            lines.append(f"   - {insight.detail}")
            lines.append(f"   - *Implication:* {insight.implication}")
            lines.append("")
        
        lines.append("## Detailed Insights")
        lines.append("")
        
        for insight in analysis.insights:
            lines.append(f"### {insight.headline}")
            lines.append("")
            lines.append(f"{insight.detail}")
            lines.append("")
            ev = insight.evidence
            if isinstance(ev, NarrativeEvidence):
                lines.append(f"- **Evidence ({ev.source_type}):** {ev.reference}")
                if ev.speaker:
                    lines.append(f"- **Speaker:** {ev.speaker}")
            else:
                lines.append(f"- **Metric:** {ev.metric} = {ev.value}")
                if ev.delta is not None:
                    lines.append(f"- **Change:** {ev.delta_percent:.1f}%")
            lines.append(f"- **Implication:** {insight.implication}")
            lines.append("")
        
        with open(file_path, 'w', encoding='utf-8') as f:
            f.write("\n".join(lines))
    
    def _generate_powerpoint(self, analysis: AnalysisResult, file_path: Path):
        """Generate PowerPoint presentation"""
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.enum.text import PP_ALIGN
            
            prs = Presentation()
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)
            
            # Title slide
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            
            title.text = "Tragaldabas Analysis"
            subtitle.text = f"{analysis.domain.value.title()} Domain\n{datetime.now().strftime('%B %Y')}"
            
            # Summary slide
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            title = slide.shapes.title
            content = slide.placeholders[1]
            
            title.text = "Executive Summary"
            tf = content.text_frame
            tf.text = "Key Findings:"

            # The Genius Move slide (Alpha Strike)
            if analysis.genius_insight and analysis.genius_insight.thesis:
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                title = slide.shapes.title
                content = slide.placeholders[1]
                title.text = "The Genius Move"
                tf = content.text_frame
                g = analysis.genius_insight
                tf.text = g.thesis
                if g.mechanism:
                    p = tf.add_paragraph()
                    p.text = g.mechanism
                    p.level = 1
                if g.estimated_upside:
                    p = tf.add_paragraph()
                    p.text = f"Upside: {g.estimated_upside}"
                    p.level = 1
                if g.kill_switch:
                    p = tf.add_paragraph()
                    p.text = f"Kill Switch: {g.kill_switch}"
                    p.level = 1
            
            for insight in analysis.insights[:5]:
                p = tf.add_paragraph()
                p.text = f"• {insight.headline}"
                p.level = 1
            
            # Insight slides
            for insight in analysis.insights:
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                title = slide.shapes.title
                content = slide.placeholders[1]
                
                title.text = insight.headline
                tf = content.text_frame
                tf.text = insight.detail
                
                ev = insight.evidence
                if isinstance(ev, NarrativeEvidence):
                    p = tf.add_paragraph()
                    p.text = f"Evidence ({ev.source_type}): {ev.reference}"
                    p.level = 1
                    if ev.speaker:
                        p = tf.add_paragraph()
                        p.text = f"Speaker: {ev.speaker}"
                        p.level = 1
                else:
                    p = tf.add_paragraph()
                    p.text = f"Metric: {ev.metric} = {ev.value}"
                    p.level = 1
                    if ev.delta is not None:
                        p = tf.add_paragraph()
                        p.text = f"Change: {ev.delta_percent:.1f}%"
                        p.level = 1
                
                p = tf.add_paragraph()
                p.text = f"Implication: {insight.implication}"
                p.level = 1
            
            prs.save(file_path)
            
        except ImportError:
            # If python-pptx not available, create a placeholder
            with open(file_path.with_suffix('.txt'), 'w') as f:
                f.write("PowerPoint generation requires python-pptx package.\n")
                f.write("Install with: pip install python-pptx\n")
                f.write(f"\nSee {file_path.with_suffix('.md')} for formatted output.")

